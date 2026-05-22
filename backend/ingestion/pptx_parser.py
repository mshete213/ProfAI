import io

from pptx import Presentation

from core.chunker import MAX_CHUNK_TOKENS, MIN_CHUNK_TOKENS, Chunk, chunk_plain_text, count_tokens


def _extract_slide_text(slide) -> tuple[str, str]:
    """Return (title, body_text_with_notes)."""
    title = ""
    body_parts: list[str] = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            text = "".join(run.text for run in para.runs).strip()
            if not text:
                continue
            if shape == slide.shapes.title and not title:
                title = text
            else:
                body_parts.append(text)

    body = "\n".join(body_parts)
    if slide.has_notes_slide:
        notes_tf = slide.notes_slide.notes_text_frame
        notes = (notes_tf.text or "").strip() if notes_tf else ""
        if notes:
            body = f"{body}\n\n[Notes]: {notes}" if body else f"[Notes]: {notes}"

    return title, body


def chunk_pptx(data: bytes, filename: str, title: str | None = None) -> list[Chunk]:
    prs = Presentation(io.BytesIO(data))
    chunks: list[Chunk] = []

    for idx, slide in enumerate(prs.slides, start=1):
        slide_title, body = _extract_slide_text(slide)
        full_text = (f"{slide_title}\n\n{body}" if slide_title else body).strip()
        if not full_text:
            continue

        tokens = count_tokens(full_text)
        metadata = {
            "filename": filename,
            "title": title or slide_title or filename,
            "slide_number": idx,
            "source_type": "pptx",
        }
        if tokens > MAX_CHUNK_TOKENS:
            for sub in chunk_plain_text(full_text, extra_metadata=metadata):
                chunks.append(sub)
        elif tokens >= MIN_CHUNK_TOKENS:
            chunks.append(Chunk(text=full_text, token_count=tokens, metadata=metadata))

    return chunks
